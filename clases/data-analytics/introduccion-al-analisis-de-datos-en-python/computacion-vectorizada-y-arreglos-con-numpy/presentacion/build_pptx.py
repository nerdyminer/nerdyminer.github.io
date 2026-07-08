"""
Genera la presentación McKinsey de NumPy en formato .pptx
Ejecutar: pipenv run python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Paleta de colores McKinsey ───────────────────────────────────────────────
NAVY       = RGBColor(0x16, 0x21, 0x3E)   # fondo secciones / cabeceras
BLUE       = RGBColor(0x0F, 0x34, 0x60)   # títulos principales
RED        = RGBColor(0xE9, 0x45, 0x60)   # acento / línea decorativa
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF8)
MID_GRAY   = RGBColor(0xCC, 0xCC, 0xDD)
DARK_GRAY  = RGBColor(0x44, 0x44, 0x55)
CODE_BG    = RGBColor(0xF0, 0xF0, 0xF8)
CODE_FG    = RGBColor(0x1A, 0x1A, 0x2E)

# ─── Dimensiones (widescreen 16:9) ────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=14, color=DARK_GRAY,
                       align=PP_ALIGN.LEFT, font_name="Calibri",
                       line_spacing=None):
    """lines: list of (text, bold, italic, size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, italic, sz, col = item, False, False, None, None
        else:
            text = item[0]
            bold   = item[1] if len(item) > 1 else False
            italic = item[2] if len(item) > 2 else False
            sz     = item[3] if len(item) > 3 else None
            col    = item[4] if len(item) > 4 else None

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = text
        run.font.size  = Pt(sz if sz else font_size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = col if col else color
        run.font.name  = font_name
    return txBox


def add_code_block(slide, code, left, top, width, height, font_size=11):
    """Renders a code block with background."""
    add_rect(slide, left, top, width, height, fill_color=CODE_BG)
    add_rect(slide, left, top, Pt(4), height, fill_color=RED)  # left accent bar
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08),
                                     width - Inches(0.2), height - Inches(0.12))
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split('\n'):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.name  = "Courier New"
        run.font.color.rgb = CODE_FG
    return txBox


def add_callout(slide, title, body, left, top, width, height,
                accent=RED, bg=RGBColor(0xFF, 0xF0, 0xF3)):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_rect(slide, left, top, Pt(5), height, fill_color=accent)
    add_text_box(slide, title,
                 left + Inches(0.15), top + Inches(0.05),
                 width - Inches(0.2), Inches(0.3),
                 font_size=12, bold=True, color=accent)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.32),
        width - Inches(0.2), height - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size  = Pt(11)
    run.font.color.rgb = DARK_GRAY
    run.font.name  = "Calibri"


def section_divider(title, subtitle=""):
    """Dark navy section-title slide."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_color=NAVY)
    # red accent bar
    add_rect(slide, Inches(1.0), Inches(3.0), Inches(0.07), Inches(1.4), fill_color=RED)
    add_text_box(slide, title,
                 Inches(1.3), Inches(3.0), Inches(10.5), Inches(0.9),
                 font_size=40, bold=True, color=WHITE, font_name="Calibri")
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(1.3), Inches(3.95), Inches(10.5), Inches(0.6),
                     font_size=18, color=MID_GRAY, font_name="Calibri")
    # footer
    add_text_box(slide, "Análisis de Datos con Python · 4° año ICM",
                 Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.35),
                 font_size=10, color=MID_GRAY, font_name="Calibri")
    return slide


def content_slide(title):
    """White content slide with McKinsey header."""
    slide = prs.slides.add_slide(BLANK)
    # header bar
    add_rect(slide, 0, 0, W, Inches(0.9), fill_color=NAVY)
    # red accent line below header
    add_rect(slide, 0, Inches(0.9), W, Pt(3), fill_color=RED)
    # title
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.1), Inches(12.3), Inches(0.72),
                 font_size=22, bold=True, color=WHITE, font_name="Calibri")
    # footer line
    add_rect(slide, 0, Inches(7.1), W, Pt(1), fill_color=MID_GRAY)
    add_text_box(slide, "Análisis de Datos con Python · 4° año ICM",
                 Inches(0.4), Inches(7.15), Inches(9), Inches(0.3),
                 font_size=9, color=MID_GRAY, font_name="Calibri")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=NAVY)
add_rect(slide, Inches(0.6), Inches(1.6), Inches(11.8), Pt(3), fill_color=RED)
add_text_box(slide, "Computación Vectorizada y Arreglos con NumPy",
             Inches(0.6), Inches(1.8), Inches(11.8), Inches(1.4),
             font_size=38, bold=True, color=WHITE, font_name="Calibri")
add_text_box(slide, "Del dato al modelo — el lenguaje del cómputo numérico en Python",
             Inches(0.6), Inches(3.3), Inches(11.8), Inches(0.6),
             font_size=20, color=RED, font_name="Calibri")
add_text_box(slide, "René Quezada C.  ·  Análisis de Datos  ·  Ingeniería Civil en Minería",
             Inches(0.6), Inches(4.0), Inches(11.8), Inches(0.4),
             font_size=14, color=MID_GRAY, font_name="Calibri")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
slide = content_slide("Agenda")
modulos = [
    ("Módulo 1 — El arreglo (ndarray)",
     "El ndarray y su geometría · Rutinas de creación · Tipos de datos · Números aleatorios · I/O · Indexación, slicing y redimensionamiento"),
    ("Módulo 2 — Operatoria vectorizada",
     "Por qué los bucles son lentos · ufuncs · Aritmética · Hadamard vs producto matricial · Álgebra lineal · Funciones matemáticas"),
    ("Módulo 3 — Broadcasting y Agregación",
     "Reglas de broadcasting · Estandarización z-score · Operaciones de agregación por eje"),
    ("Módulo 4 — Filtrado, Selección y Ordenamiento",
     "Comparación booleana · Masking · Fancy indexing · Ordenamiento · Reemplazo condicional"),
    ("Módulo 5 — Modelamiento predictivo",
     "Regresión lineal con NumPy · Gradiente descendente · Regresión logística binaria · Evaluación de modelos"),
]
y_pos = Inches(1.05)
for i, (mod, desc) in enumerate(modulos):
    # colored index box
    add_rect(slide, Inches(0.4), y_pos, Inches(0.35), Inches(0.9),
             fill_color=NAVY if i % 2 == 0 else BLUE)
    add_text_box(slide, str(i+1),
                 Inches(0.4), y_pos + Inches(0.18), Inches(0.35), Inches(0.5),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, mod,
                 Inches(0.9), y_pos, Inches(11.8), Inches(0.32),
                 font_size=13, bold=True, color=BLUE, font_name="Calibri")
    add_text_box(slide, desc,
                 Inches(0.9), y_pos + Inches(0.32), Inches(11.8), Inches(0.55),
                 font_size=11, color=DARK_GRAY, font_name="Calibri")
    y_pos += Inches(1.05)


# ══════════════════════════════════════════════════════════════════════════════
# MÓDULO 1
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Módulo 1", "El arreglo · Geometría · Creación · Tipos · Indexación")

# SLIDE 3 — ¿Qué es NumPy?
slide = content_slide("NumPy es la base de todo el ecosistema científico de Python")
add_callout(slide, "Idea central",
            "NumPy (Numerical Python) introduce el ndarray: una estructura para representar vectores, matrices y tensores, "
            "ejecutar operaciones vectorizadas con gran eficiencia y servir de fundamento al ecosistema completo de "
            "ciencia de datos en Python.",
            Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.85))

libs = [
    ("Pandas",       "Análisis tabular"),
    ("Matplotlib",   "Visualización"),
    ("SciPy",        "Computación científica"),
    ("Scikit-Learn", "Machine Learning"),
    ("PyTorch",      "Deep Learning"),
]
add_text_box(slide, "Librerías construidas sobre NumPy:",
             Inches(0.4), Inches(2.0), Inches(5.5), Inches(0.35),
             font_size=13, bold=True, color=BLUE)
y = Inches(2.4)
for lib, desc in libs:
    add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.28), fill_color=RED)
    add_text_box(slide, f"{lib}  —  {desc}",
                 Inches(0.6), y, Inches(5.2), Inches(0.3),
                 font_size=12, color=DARK_GRAY)
    y += Inches(0.38)

add_text_box(slide, "Instalación e importación:",
             Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.35),
             font_size=13, bold=True, color=BLUE)
add_code_block(slide, "pip install numpy\n\nimport numpy as np\n# El alias 'np' es convención universal",
               Inches(7.0), Inches(2.4), Inches(5.8), Inches(1.3), font_size=13)
add_text_box(slide,
             "Entender NumPy no es aprender una librería. Es construir el vocabulario\n"
             "mínimo del análisis moderno de datos en Python.",
             Inches(0.4), Inches(6.3), Inches(12.2), Inches(0.65),
             font_size=12, italic=True, color=DARK_GRAY)

# SLIDE 4 — El ndarray
slide = content_slide("El ndarray: una grilla de valores con geometría bien definida")
add_text_box(slide, "El arreglo es el elemento central de NumPy:",
             Inches(0.4), Inches(1.05), Inches(5.8), Inches(0.35),
             font_size=13, bold=True, color=BLUE)
puntos = [
    "Un único tipo de dato (todos float64, o todos int32, etc.)",
    "Una geometría determinada por sus ejes",
    "Soporte nativo para operaciones vectorizadas",
]
y = Inches(1.45)
for p in puntos:
    add_rect(slide, Inches(0.4), y + Inches(0.07), Inches(0.12), Inches(0.12), fill_color=RED)
    add_text_box(slide, p, Inches(0.65), y, Inches(5.6), Inches(0.3), font_size=12, color=DARK_GRAY)
    y += Inches(0.36)

# tabla objetos matemáticos
add_text_box(slide, "Los arreglos representan objetos matemáticos familiares:",
             Inches(0.4), Inches(2.65), Inches(5.8), Inches(0.35),
             font_size=13, bold=True, color=BLUE)
tabla = [("1D", "Vector",       "(n,)"),
         ("2D", "Matriz",       "(n, d)"),
         ("3D", "Tensor ord. 3","(k, n, d)")]
ys = Inches(3.05)
add_rect(slide, Inches(0.4), ys, Inches(5.5), Inches(0.32), fill_color=NAVY)
for lbl, col1, col2 in [("Dims", "Objeto", "Geometría")]:
    for col, x in [(lbl, 0.4), (col1, 1.5), (col2, 3.8)]:
        add_text_box(slide, col, Inches(x), ys + Inches(0.05),
                     Inches(1.5), Inches(0.25), font_size=11, bold=True, color=WHITE)
ys += Inches(0.32)
for i, (d, obj, geo) in enumerate(tabla):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), ys, Inches(5.5), Inches(0.32), fill_color=bg)
    add_text_box(slide, d,   Inches(0.5),  ys + Inches(0.07), Inches(0.9), Inches(0.25), font_size=11, color=DARK_GRAY)
    add_text_box(slide, obj, Inches(1.5),  ys + Inches(0.07), Inches(2.0), Inches(0.25), font_size=11, color=DARK_GRAY)
    add_text_box(slide, geo, Inches(3.8),  ys + Inches(0.07), Inches(2.1), Inches(0.25), font_size=11, color=BLUE, bold=True)
    ys += Inches(0.32)

add_code_block(slide,
    "x = np.array([-1, 1, 5, -8, 2])        # Vector  shape: (5,)\n"
    "u = np.array([[-1, 1, 5, -8, 2]])       # M. fila shape: (1, 5)\n"
    "w = np.array([[-1],[ 1],[ 5],[-8],[2]]) # M. col  shape: (5, 1)\n"
    "\n"
    "A = np.array([                           # Matriz  shape: (4, 4)\n"
    "    [-2,  1, -7,  6],\n"
    "    [ 1,  3,  1, -4],\n"
    "])",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.6), font_size=12)

add_callout(slide, "¡Atención!",
            "Un arreglo 1D (5,) NO es una matriz fila (1, 5). Son objetos distintos con reglas distintas.",
            Inches(6.6), Inches(3.8), Inches(6.3), Inches(0.75),
            accent=RGBColor(0xE6, 0x7E, 0x22), bg=RGBColor(0xFE, 0xF9, 0xF0))

add_callout(slide, "Restricción clave",
            "Un arreglo solo puede contener datos de UN único tipo. No se permiten mezclas.",
            Inches(0.4), Inches(4.45), Inches(5.8), Inches(0.75))

# SLIDE 5 — Geometría y ejes
slide = content_slide("Los ejes de un arreglo: el sistema de referencia interno de NumPy")
add_callout(slide, "Idea clave",
            "Los ejes (axis) son las direcciones a lo largo de las cuales NumPy opera. "
            "El eje 0 recorre las filas; el eje 1 recorre las columnas. "
            "Toda operación de agregación depende de este concepto.",
            Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.85))

# diagrama ASCII de ejes
add_rect(slide, Inches(0.4), Inches(2.1), Inches(5.8), Inches(2.8), fill_color=LIGHT_GRAY)
diag = (
    "               axis=1 →\n"
    "         ┌──────────────┐\n"
    "  axis=0 │  a   b   c   │\n"
    "    ↓    │  d   e   f   │\n"
    "         │  g   h   i   │\n"
    "         └──────────────┘\n"
    "\n"
    " A.sum(axis=1) → suma por filas\n"
    " A.sum(axis=0) → suma por columnas"
)
add_text_box(slide, diag,
             Inches(0.55), Inches(2.15), Inches(5.5), Inches(2.65),
             font_size=12, font_name="Courier New", color=CODE_FG)

add_code_block(slide,
    "A = np.array([\n"
    "    [-2, 1, -7, 6],\n"
    "    [ 1, 3,  1,-4],\n"
    "    [-5,-5,  0, 4],\n"
    "    [-9, 2, -8, 9]\n"
    "])\n\n"
    "A.sum(axis=1)   # suma por filas → [-2, 1, -6, -6]\n"
    "A.sum(axis=0)   # suma por cols  → [-15, 1, -14, 15]\n\n"
    "# Concatenar dos arreglos compatibles\n"
    "B = np.array([[-1,4,5,-8],[0,-5,6,-9]])\n"
    "np.concatenate([A, B], axis=0)  # apila filas",
    Inches(6.6), Inches(2.1), Inches(6.3), Inches(2.85), font_size=11)

add_callout(slide, "Regla práctica",
            "Forma de pensar axis: 'aplastar el arreglo en la dirección del eje especificado'. "
            "axis=0 colapsa filas → resultado por columna. axis=1 colapsa columnas → resultado por fila.",
            Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.85),
            accent=RGBColor(0x1D, 0x9A, 0x6C), bg=RGBColor(0xF0, 0xFA, 0xF5))

# SLIDE 6 — Atributos y tipos de datos
slide = content_slide("Atributos esenciales y tipos de datos de un arreglo")
add_text_box(slide, "Atributos del ndarray:",
             Inches(0.4), Inches(1.05), Inches(5.8), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
attrs = [
    (".ndim",     "Número de dimensiones"),
    (".shape",    "Geometría: (filas, cols, ...)"),
    (".size",     "Total de elementos"),
    (".dtype",    "Tipo de dato"),
    (".itemsize", "Bytes por elemento"),
    (".nbytes",   "Bytes totales = size × itemsize"),
]
add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.32), fill_color=NAVY)
add_text_box(slide, "Atributo",  Inches(0.5), Inches(1.47), Inches(1.5), Inches(0.22), font_size=11, bold=True, color=WHITE)
add_text_box(slide, "Descripción", Inches(2.1), Inches(1.47), Inches(3.9), Inches(0.22), font_size=11, bold=True, color=WHITE)
y = Inches(1.72)
for i, (a, d) in enumerate(attrs):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(5.8), Inches(0.3), fill_color=bg)
    add_text_box(slide, a, Inches(0.5), y + Inches(0.05), Inches(1.5), Inches(0.25), font_size=11, bold=True, color=BLUE, font_name="Courier New")
    add_text_box(slide, d, Inches(2.1), y + Inches(0.05), Inches(3.9), Inches(0.25), font_size=11, color=DARK_GRAY)
    y += Inches(0.3)

add_text_box(slide, "Tipos de datos clave (dtype):",
             Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
tipos = [
    ("bool",       "True / False"),
    ("int8/16/32/64", "Entero con signo"),
    ("uint8/16/32/64","Entero sin signo (≥0)"),
    ("float32",    "Precisión simple (~7 dígitos)"),
    ("float64",    "Precisión doble (~15 dígitos) ✓"),
    ("complex128", "Dos float64 (números complejos)"),
]
add_rect(slide, Inches(6.6), Inches(1.4), Inches(6.3), Inches(0.32), fill_color=NAVY)
add_text_box(slide, "Tipo",  Inches(6.7), Inches(1.47), Inches(2.3), Inches(0.22), font_size=11, bold=True, color=WHITE)
add_text_box(slide, "Descripción", Inches(9.2), Inches(1.47), Inches(3.5), Inches(0.22), font_size=11, bold=True, color=WHITE)
y = Inches(1.72)
for i, (t, d) in enumerate(tipos):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(6.6), y, Inches(6.3), Inches(0.3), fill_color=bg)
    add_text_box(slide, t, Inches(6.7), y + Inches(0.05), Inches(2.3), Inches(0.25), font_size=11, bold=True, color=BLUE, font_name="Courier New")
    add_text_box(slide, d, Inches(9.2), y + Inches(0.05), Inches(3.5), Inches(0.25), font_size=11, color=DARK_GRAY)
    y += Inches(0.3)

add_code_block(slide,
    "x2 = rng.integers(0, 10, size=(3, 4))\n"
    "x2.ndim      # 2\n"
    "x2.shape     # (3, 4)\n"
    "x2.size      # 12\n"
    "x2.dtype     # int64\n"
    "x2.itemsize  # 8 bytes\n"
    "x2.nbytes    # 96 (= 12 × 8)",
    Inches(0.4), Inches(4.12), Inches(5.8), Inches(2.0), font_size=12)

add_callout(slide, "¡Atención!",
            "Si insertas un float en un arreglo int, NumPy trunca silenciosamente el valor. "
            "Ejemplo: x[2] = 3.99 guarda 3, no 4.",
            Inches(6.6), Inches(4.12), Inches(6.3), Inches(0.8),
            accent=RGBColor(0xE6, 0x7E, 0x22), bg=RGBColor(0xFE, 0xF9, 0xF0))

# SLIDE 7 — Rutinas de creación
slide = content_slide("Rutinas de creación: construir arreglos sin ingresar cada valor")
add_code_block(slide,
    "# Arreglos constantes\n"
    "np.zeros(8)                      # [0. 0. 0. 0. 0. 0. 0. 0.]\n"
    "np.zeros((5, 4))                 # Matriz 5×4 de ceros\n"
    "np.ones((4, 6), dtype=int)       # Matriz 4×6 de unos enteros\n"
    "np.full((5, 5), fill_value=9)    # Matriz 5×5 llena con 9",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(1.55), font_size=11)

add_code_block(slide,
    "# Matrices diagonales\n"
    "np.eye(N=5, M=6, k=0)   # diagonal reducida (1s en diag principal)\n"
    "np.eye(N=4, M=8, k=1)   # k>0: diagonal superior\n"
    "np.identity(n=6)         # matriz identidad cuadrada",
    Inches(0.4), Inches(2.7), Inches(6.0), Inches(1.4), font_size=11)

add_code_block(slide,
    "# Rangos\n"
    "np.arange(start=1, stop=20, step=2)     # [1,3,5,...,19]\n"
    "np.arange(start=100, stop=0, step=-10)  # [100,90,...,10]\n\n"
    "# Intervalo con número de puntos fijo\n"
    "np.linspace(start=0, stop=1, num=5)     # [0, .25, .5, .75, 1.0]",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(1.85), font_size=11)

add_code_block(slide,
    "# Grillas para evaluar f(x,y)\n"
    "x = np.linspace(-3, 3, 100)\n"
    "y = np.linspace(-3, 3, 100)\n"
    "X, Y = np.meshgrid(x, y, indexing='ij')\n"
    "# X, Y: shape (100, 100)\n\n"
    "Z = np.exp(-(X**2 + Y**2))  # evaluación vectorizada",
    Inches(6.6), Inches(3.0), Inches(6.3), Inches(2.0), font_size=11)

add_callout(slide, "arange vs linspace",
            "arange NO incluye el valor stop. linspace SÍ lo incluye. "
            "Error frecuente al construir rangos.",
            Inches(0.4), Inches(4.2), Inches(5.9), Inches(0.75),
            accent=RGBColor(0xE6, 0x7E, 0x22), bg=RGBColor(0xFE, 0xF9, 0xF0))

# SLIDE 8 — Números aleatorios
slide = content_slide("Generación de números pseudoaleatorios: reproducibilidad por semilla")
add_code_block(slide,
    "rng = np.random.default_rng(seed=42)  # semilla fija\n\n"
    "rng.random(size=(3,3))           # Uniforme en [0, 1)\n"
    "rng.normal(loc=0, scale=1, size=(3,3))  # Normal N(0,1)\n"
    "rng.integers(low=0, high=10, size=(4,4))  # Enteros\n"
    "rng.uniform(size=(30, 2))        # Uniforme continua\n\n"
    "# Misma semilla → mismo resultado siempre.",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.3), font_size=12)

add_code_block(slide,
    "# Distribuciones con aplicación en minería\n\n"
    "# Poisson: conteo de fallas por mes\n"
    "rng.poisson(lam=2.5, size=(4,4))\n\n"
    "# Beta: recuperaciones y proporciones\n"
    "# También usada en tiempos PERT (ruta crítica)\n"
    "rng.beta(a=2.0, b=5.0, size=(4,4))\n\n"
    "# Normal multivariante: leyes correlacionadas\n"
    "mean = [0.45, 0.12]          # Cu%, Mo%\n"
    "cov  = [[0.01, 0.002],[0.002, 0.001]]\n"
    "rng.multivariate_normal(mean, cov, size=100)",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(3.3), font_size=11)

add_callout(slide, "Reproducibilidad",
            "Misma semilla → mismo arreglo en toda ejecución, en toda máquina. "
            "Indispensable para que los experimentos sean repetibles y auditables.",
            Inches(0.4), Inches(3.5), Inches(5.9), Inches(0.85))

# SLIDE 9 — Indexación y slicing
slide = content_slide("Indexación y slicing: acceso preciso a elementos y sub-arreglos")
add_code_block(slide,
    "x = np.array([5, 2, 8, 1, 9, 4])\n"
    "x[0]   # → 5   (primer elemento)\n"
    "x[-1]  # → 4   (último)\n\n"
    "# Arreglos 2D: (fila, columna)\n"
    "x2[0, 1]   # fila 0, col 1\n"
    "x2[-1, -2] # última fila, penúltima col\n\n"
    "# Modificar en sitio\n"
    "x2[1, 0] = 99",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.2), font_size=12)

add_code_block(slide,
    "# Slicing: x[inicio:final:paso]\n"
    "x = np.arange(10)  # [0,1,2,...,9]\n\n"
    "x[:5]      # → [0,1,2,3,4]\n"
    "x[5:]      # → [5,6,7,8,9]\n"
    "x[4:7]     # → [4,5,6]\n"
    "x[::2]     # → [0,2,4,6,8]\n"
    "x[::-1]    # → [9,8,...,0]   (invertir)\n\n"
    "# En 2D:\n"
    "A[:2, :3]  # filas 0-1, columnas 0-2\n"
    "A[:, 1]    # TODA la columna 1\n"
    "A[1, :]    # TODA la fila 1",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.9), font_size=11)

add_callout(slide, "Slices son VISTAS, no copias",
            "Modificar un sub-arreglo modifica el original. "
            "Para una copia independiente: sub = A[:2, :2].copy()\n"
            "Razón: trabajar con fragmentos de datasets enormes sin duplicar memoria.",
            Inches(0.4), Inches(3.4), Inches(5.9), Inches(1.0))

add_callout(slide, "Redimensionamiento",
            "A.T → transpuesta\n"
            "A.ravel() → aplana a 1D\n"
            "A.reshape(m, n) → nueva geometría (conserva size)\n"
            "np.vstack / np.hstack / np.concatenate → unión de arreglos",
            Inches(6.6), Inches(4.1), Inches(6.3), Inches(1.05),
            accent=RGBColor(0x1D, 0x9A, 0x6C), bg=RGBColor(0xF0, 0xFA, 0xF5))

# SLIDE 10 — I/O
slide = content_slide("I/O: cargar y guardar arreglos desde y hacia disco")
add_code_block(slide,
    "# Lectura simple (archivo limpio)\n"
    "A = np.loadtxt('datos.csv', delimiter=',')\n\n"
    "# Lectura flexible (tolera NaN, cabeceras)\n"
    "X = np.genfromtxt(\n"
    "    'mediciones.csv',\n"
    "    delimiter=',',\n"
    "    skip_header=1,\n"
    "    filling_values=-999   # NaN → -999\n"
    ")\n\n"
    "# Leer sólo ciertas columnas\n"
    "X = np.loadtxt('datos.csv', delimiter=',',\n"
    "               usecols=(0, 2))",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(3.0), font_size=11)

add_code_block(slide,
    "# Escritura como texto (CSV)\n"
    "np.savetxt('salida.csv', A, delimiter=',',\n"
    "           fmt='%.2f')   # 2 decimales\n\n"
    "# Formato binario nativo (.npy)\n"
    "# Preserva geometría y dtype exactamente\n"
    "np.save('mi_arreglo.npy', A)\n"
    "B = np.load('mi_arreglo.npy')  # recupera A",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.2), font_size=11)

reglas = [
    ("loadtxt",       "Datos limpios y regulares, sin cabeceras"),
    ("genfromtxt",    "Cabeceras, NaN, o estructura irregular"),
    ("savetxt",       "Guardar en texto legible (CSV)"),
    ("save / load",   "Resultados intermedios — formato binario"),
    ("Pandas",        "Datos con tipos mixtos, nombres de columna"),
]
add_text_box(slide, "Regla de uso:",
             Inches(6.6), Inches(3.4), Inches(6.3), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
y = Inches(3.75)
for fn, desc in reglas:
    add_rect(slide, Inches(6.6), y, Inches(0.06), Inches(0.26), fill_color=RED)
    add_text_box(slide, f"{fn:<16} — {desc}",
                 Inches(6.8), y, Inches(6.1), Inches(0.28),
                 font_size=11, font_name="Courier New", color=DARK_GRAY)
    y += Inches(0.33)


# ══════════════════════════════════════════════════════════════════════════════
# MÓDULO 2
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Módulo 2", "Operatoria Vectorizada · ufuncs · Aritmética · Álgebra lineal · Funciones matemáticas")

# SLIDE 11 — Por qué los bucles son lentos
slide = content_slide("Los bucles en Python son lentos: CPython revisa tipos en cada iteración")
add_code_block(slide,
    "# Con bucle: lento\n"
    "def compute_neg_exp(values):\n"
    "    output = np.empty(len(values))\n"
    "    for i in range(len(values)):\n"
    "        output[i] = 2**(-values[i])   # chequeo de tipo en cada ciclo\n"
    "    return output\n\n"
    "huge = rng.normal(5, 2, size=1_000_000)\n"
    "%timeit compute_neg_exp(huge)\n"
    "# → ~500 ms por ejecución",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.6), font_size=11)

add_code_block(slide,
    "# Vectorizado: ~150x más rápido\n"
    "%timeit 2**(-huge)\n"
    "# → ~3 ms por ejecución\n\n"
    "# La operación se delega a rutinas\n"
    "# compiladas en C que evitan el\n"
    "# despacho dinámico de Python.",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.1), font_size=11)

add_callout(slide, "Regla de oro",
            "Si puedes escribirlo como una operación sobre un arreglo completo, NUNCA uses un bucle explícito.",
            Inches(0.4), Inches(3.8), Inches(12.5), Inches(0.75))

add_text_box(slide, "Ganancia de rendimiento:",
             Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
comparacion = [
    ("Bucle Python", "~500 ms", "Revisión de tipo en cada elemento"),
    ("NumPy ufunc",  "~3 ms",   "Rutina C precompilada, sin overhead"),
]
y = Inches(5.05)
add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.32), fill_color=NAVY)
for lbl, tw, pos in [("Método", "Tiempo (1M elem.)", "Por qué")]:
    add_text_box(slide, lbl, Inches(0.5), y + Inches(0.05), Inches(2.5), Inches(0.25), font_size=11, bold=True, color=WHITE)
    add_text_box(slide, tw,  Inches(3.2), y + Inches(0.05), Inches(2.5), Inches(0.25), font_size=11, bold=True, color=WHITE)
    add_text_box(slide, pos, Inches(6.0), y + Inches(0.05), Inches(6.7), Inches(0.25), font_size=11, bold=True, color=WHITE)
y += Inches(0.32)
for i, (m, t, why) in enumerate(comparacion):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.34), fill_color=bg)
    add_text_box(slide, m,   Inches(0.5), y + Inches(0.06), Inches(2.5), Inches(0.28), font_size=11, color=DARK_GRAY)
    tc = RED if i == 0 else RGBColor(0x1D, 0x9A, 0x6C)
    add_text_box(slide, t,   Inches(3.2), y + Inches(0.06), Inches(2.5), Inches(0.28), font_size=12, bold=True, color=tc)
    add_text_box(slide, why, Inches(6.0), y + Inches(0.06), Inches(6.7), Inches(0.28), font_size=11, color=DARK_GRAY)
    y += Inches(0.34)

# SLIDE 12 — ufuncs y aritmética
slide = content_slide("ufuncs: las operaciones vectorizadas son siempre componente a componente")
add_code_block(slide,
    "x = np.array([-1, 0, 9, 5, -4, 5])\n\n"
    "x + 5      # → [ 4,  5, 14, 10,  1, 10]\n"
    "x - 5      # → [-6, -5,  4,  0, -9,  0]\n"
    "x * 2      # → [-2,  0, 18, 10, -8, 10]\n"
    "x / 2      # → [-0.5, 0, 4.5, 2.5, -2, 2.5]\n"
    "x // 2     # → [-1,  0,  4,  2, -2,  2]  (cociente)\n"
    "x % 2      # → [ 1,  0,  1,  1,  0,  1]  (módulo)\n"
    "-x         # → [ 1,  0, -9, -5,  4, -5]  (negación)\n"
    "x ** 2     # → [ 1,  0, 81, 25, 16, 25]",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.9), font_size=11)

add_code_block(slide,
    "a = np.array([0, -1, 8, -2, 5])\n"
    "b = np.array([-6, -9, 0, 1, -2])\n\n"
    "a + b   # → [-6, -10, 8, -1, 3]   (suma vectorial)\n"
    "2 * a   # → [0, -2, 16, -4, 10]   (producto por escalar)\n\n"
    "# Mezcla de ufuncs (respeta jerarquía aritmética)\n"
    "-(0.5*x + 1) ** 2",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.1), font_size=11)

add_callout(slide, "Fundamento algebraico",
            "La suma y el producto por escalar de arreglos replican exactamente las operaciones del "
            "espacio vectorial ℝⁿ (o ℝⁿˣᵐ para matrices). "
            "El cálculo es siempre componente a componente.",
            Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.95),
            accent=RGBColor(0x1D, 0x9A, 0x6C), bg=RGBColor(0xF0, 0xFA, 0xF5))

# SLIDE 13 — Hadamard vs matricial
slide = content_slide("Producto de Hadamard ≠ Producto matricial: distinción crítica")
add_callout(slide, "Error frecuente — leer antes de continuar",
            "El operador * en NumPy es SIEMPRE el producto de Hadamard (componente a componente), "
            "NO el producto matricial. Confundir ambos produce errores silenciosos difíciles de detectar.",
            Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.85),
            accent=RED, bg=RGBColor(0xFF, 0xF0, 0xF3))

add_text_box(slide, "Producto de Hadamard  A ⊙ B  (componente a componente):",
             Inches(0.4), Inches(2.05), Inches(5.9), Inches(0.32),
             font_size=12, bold=True, color=BLUE)
add_code_block(slide,
    "A = np.array([[1, 2], [3, 4]])\n"
    "B = np.array([[5, 6], [7, 8]])\n\n"
    "A * B          # → [[5, 12], [21, 32]]\n"
    "               #    componente a componente\n\n"
    "# Producto interno de vectores: Σ aᵢbᵢ → escalar\n"
    "a, b = np.array([1,2,3]), np.array([4,5,6])\n"
    "np.dot(a, b)   # → 32  (escalar)",
    Inches(0.4), Inches(2.42), Inches(5.9), Inches(2.4), font_size=11)

add_text_box(slide, "Producto matricial  C = AB  (estándar del álgebra lineal):",
             Inches(6.6), Inches(2.05), Inches(6.3), Inches(0.32),
             font_size=12, bold=True, color=BLUE)
add_code_block(slide,
    "A = np.array([[1, 2, 3], [3, 2, 1]])  # (2×3)\n"
    "B = np.array([[0, 2], [1,-1], [0, 1]]) # (3×2)\n\n"
    "# Tres formas equivalentes:\n"
    "np.matmul(A, B)   # → (2×2)\n"
    "A @ B             # → (2×2)  ← PREFERIDO\n"
    "np.dot(A, B)      # → (2×2)  (para 2D)\n\n"
    "# ¡No conmutativo!\n"
    "A @ B  !=  B @ A  # en general",
    Inches(6.6), Inches(2.42), Inches(6.3), Inches(2.65), font_size=11)

resumen = [
    ("A * B",          "Hadamard (componente a componente)     — requiere misma geometría"),
    ("A @ B",          "Producto matricial estándar            — requiere cols(A) = filas(B)"),
    ("np.dot(a, b)",   "Producto interno de vectores 1D        — devuelve escalar"),
]
y = Inches(5.1)
for op, desc in resumen:
    add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.3), fill_color=RED)
    add_text_box(slide, f"{op:<22} {desc}",
                 Inches(0.6), y + Inches(0.04), Inches(12.2), Inches(0.28),
                 font_size=11, font_name="Courier New", color=DARK_GRAY)
    y += Inches(0.38)

# SLIDE 14 — Norma y álgebra lineal
slide = content_slide("Norma, valor absoluto y submódulo de álgebra lineal np.linalg")
add_code_block(slide,
    "# Valor absoluto (componente a componente)\n"
    "A = np.array([[-1,-6,7],[5,-3,-4]])\n"
    "np.abs(A)          # → [[1,6,7],[5,3,4]]\n\n"
    "# Con complejos: calcula módulo |z| = sqrt(a²+b²)\n"
    "Z = np.array([2-3j, 4+1j])\n"
    "np.abs(Z)          # → [sqrt(13), sqrt(17)]\n\n"
    "# Norma Euclidiana (L2)\n"
    "x = np.array([1, -1, 4, -8, 0, -2])\n"
    "np.linalg.norm(x)  # = sqrt(1+1+16+64+0+4) ≈ 9.27",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.95), font_size=11)

add_text_box(slide, "El submódulo np.linalg:",
             Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
linalg_fns = [
    ("np.linalg.norm(x)",       "Norma de vector o matriz"),
    ("np.linalg.det(A)",        "Determinante"),
    ("np.linalg.eig(A)",        "Valores y vectores propios"),
    ("np.linalg.svd(A)",        "Descomposición SVD"),
    ("np.linalg.solve(A, b)",   "Resolver sistema Ax = b"),
    ("np.linalg.pinv(A)",       "Pseudoinversa (mínimos cuadrados)"),
    ("np.linalg.inv(A)",        "Inversa (solo si A es cuadrada y regular)"),
]
y = Inches(1.42)
for i, (fn, desc) in enumerate(linalg_fns):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(6.6), y, Inches(6.3), Inches(0.3), fill_color=bg)
    add_text_box(slide, fn,   Inches(6.7),  y + Inches(0.05), Inches(3.0), Inches(0.24), font_size=10, bold=True, color=BLUE, font_name="Courier New")
    add_text_box(slide, desc, Inches(9.85), y + Inches(0.05), Inches(2.9), Inches(0.24), font_size=10, color=DARK_GRAY)
    y += Inches(0.3)

add_callout(slide, "Aplicación en minería",
            "np.linalg.solve() permite resolver sistemas de balance de masa y restricciones de proceso. "
            "np.linalg.pinv() implementa ajuste por mínimos cuadrados en estimación de parámetros.",
            Inches(0.4), Inches(4.1), Inches(5.9), Inches(1.1),
            accent=RGBColor(0x1D, 0x9A, 0x6C), bg=RGBColor(0xF0, 0xFA, 0xF5))

# SLIDE 15 — Funciones matemáticas
slide = content_slide("Funciones matemáticas: trigonometría, logaritmos, exponenciales")
add_code_block(slide,
    "# Funciones trigonométricas\n"
    "theta = np.linspace(0, np.pi, 4)\n\n"
    "np.sin(theta)   # seno\n"
    "np.cos(theta)   # coseno\n"
    "np.tan(theta)   # tangente\n\n"
    "np.arcsin([-1, 0, 1])  # inversa seno\n"
    "np.arccos([-1, 0, 1])  # inversa coseno\n"
    "np.arctan([-1, 0, 1])  # inversa tangente\n\n"
    "np.rad2deg(theta)  # radianes → grados\n"
    "np.deg2rad(theta)  # grados → radianes",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(3.0), font_size=11)

add_code_block(slide,
    "# Exponenciales y logaritmos\n"
    "x = [1, 2, 3, 4, 5]\n\n"
    "np.exp(x)         # e^x\n"
    "np.exp2(x)        # 2^x\n"
    "np.power(3, x)    # 3^x\n"
    "np.sqrt(x)        # sqrt(x)\n\n"
    "np.log(x)         # ln(x)   base natural\n"
    "np.log2(x)        # log₂(x)\n"
    "np.log10(x)       # log₁₀(x)\n\n"
    "# Versiones numéricamente estables (x ≈ 0)\n"
    "np.expm1(x)       # exp(x) - 1\n"
    "np.log1p(x)       # ln(1 + x)",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(3.5), font_size=11)

add_callout(slide, "Precisión numérica",
            "Los resultados son aproximaciones de punto flotante. "
            "np.sin(np.pi) no es exactamente 0 sino ~1.2e-16. "
            "expm1 y log1p evitan pérdida de precisión cuando el argumento está cerca de cero.",
            Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.95),
            accent=RGBColor(0xE6, 0x7E, 0x22), bg=RGBColor(0xFE, 0xF9, 0xF0))


# ══════════════════════════════════════════════════════════════════════════════
# MÓDULO 3
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Módulo 3", "Broadcasting · Agregación · Patrones vectorizados")

# SLIDE 16 — Broadcasting
slide = content_slide("Broadcasting: operar con arreglos de distinta geometría sin copias manuales")
add_callout(slide, "Definición",
            "Broadcasting es el mecanismo por el cual NumPy puede operar con arreglos de distinta forma, "
            "extendiendo virtualmente las dimensiones de tamaño 1. "
            "No hay copia real en memoria — la eficiencia se preserva.",
            Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.85))

add_code_block(slide,
    "# Caso 1: arreglo + escalar\n"
    "a = np.array([1.0, 3.0, 5.0, 7.0])\n"
    "a + 2.0   # el escalar 2.0 se 'extiende'\n"
    "          # → [3.0, 5.0, 7.0, 9.0]\n\n"
    "# Caso 2: matriz + vector 1D\n"
    "M = np.eye(4)            # (4, 4)\n"
    "a = np.array([1,2,3,4])  # (4,)\n"
    "M + a   # a se extiende sobre cada fila de M",
    Inches(0.4), Inches(2.05), Inches(5.9), Inches(2.6), font_size=11)

add_code_block(slide,
    "# Caso 3: columna + fila (outer operation)\n"
    "u = np.ones(shape=(1, 3))  # (1, 3)\n"
    "v = np.ones(shape=(3, 1))  # (3, 1)\n"
    "u + v   # resultado: (3, 3) — ambos se extienden\n\n"
    "# Caso 4: grilla sin meshgrid\n"
    "X = np.linspace(-1,1,100).reshape(-1, 1)  # (100,1)\n"
    "Y = np.linspace(-1,1,100).reshape(1, -1)  # (1,100)\n"
    "Z = np.arcsinh(X**2 + Y**2)  # (100,100)",
    Inches(6.6), Inches(2.05), Inches(6.3), Inches(2.6), font_size=11)

# SLIDE 17 — Reglas de broadcasting
slide = content_slide("Las tres reglas del broadcasting: leer de derecha a izquierda")
add_text_box(slide,
    "NumPy compara geometrías dimensión a dimensión, desde la derecha. Dos dimensiones son compatibles si:",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.42),
    font_size=13, color=DARK_GRAY)

reglas = [
    ("1", "Son iguales"),
    ("2", "Una de ellas es igual a 1 (esa se 'extiende' al valor de la otra)"),
]
y = Inches(1.52)
for n, r in reglas:
    add_rect(slide, Inches(0.4), y, Inches(0.35), Inches(0.35), fill_color=NAVY)
    add_text_box(slide, n, Inches(0.4), y + Inches(0.04), Inches(0.35), Inches(0.28),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, r, Inches(0.88), y + Inches(0.06), Inches(11.8), Inches(0.3),
                 font_size=13, color=DARK_GRAY)
    y += Inches(0.5)

add_text_box(slide, "Procedimiento mental: (1) igualar dims agregando 1s a la izquierda · (2) comparar de derecha a izquierda · (3) extender solo dims compatibles",
             Inches(0.4), Inches(2.62), Inches(12.5), Inches(0.4),
             font_size=11, italic=True, color=DARK_GRAY)

ejemplos = [
    ("✓", "(2, 3)  +  (3,) → (1,3)   ultima: 3=3 ✓  | primera: 2 y 1 ✓   → (2,3)",  True),
    ("✓", "(3, 1)  +  (3,) → (1,3)   ultima: 1 y 3 ✓ | primera: 3 y 1 ✓   → (3,3)",  True),
    ("✗", "(3, 2)  +  (3,) → (1,3)   ultima: 2 ≠ 3 y ninguna es 1  → ERROR",  False),
]
y = Inches(3.1)
for mark, txt, ok in ejemplos:
    col = RGBColor(0x1D, 0x9A, 0x6C) if ok else RED
    add_rect(slide, Inches(0.4), y, Inches(0.38), Inches(0.38), fill_color=col)
    add_text_box(slide, mark, Inches(0.4), y + Inches(0.04), Inches(0.38), Inches(0.3),
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, txt, Inches(0.9), y + Inches(0.06), Inches(12.0), Inches(0.3),
                 font_size=11, font_name="Courier New", color=DARK_GRAY)
    y += Inches(0.52)

add_code_block(slide,
    "# Estandarización z-score — Broadcasting en acción\n"
    "# T: (10, 3) — 10 muestras de proceso minero × 3 variables\n"
    "T = np.vstack([t, p, g]).T\n\n"
    "# mean y std son (3,) — broadcasting extiende a (10, 3)\n"
    "Z = (T - T.mean(axis=0)) / T.std(axis=0)\n\n"
    "Z.mean(axis=0)  # ≈ [0, 0, 0]\n"
    "Z.std(axis=0)   # ≈ [1, 1, 1]",
    Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.1), font_size=11)

# SLIDE 18 — Agregación
slide = content_slide("Operaciones de agregación: resumir información por eje")
add_code_block(slide,
    "A = rng.integers(1, 10, size=(5, 10))\n\n"
    "# Suma\n"
    "A.sum()         # escalar: suma total\n"
    "A.sum(axis=1)   # (5,): una suma por fila\n"
    "A.sum(axis=0)   # (10,): una suma por columna\n\n"
    "# Suma acumulada (no colapsa el eje)\n"
    "A.cumsum(axis=1)   # misma shape que A\n\n"
    "# Estadísticas\n"
    "A.mean(axis=0)     # media por columna\n"
    "A.std(axis=0)      # desv.est. por columna\n"
    "A.max(); A.argmax()  # máximo y su posición\n"
    "A.min(); A.argmin(axis=1)",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(3.6), font_size=11)

add_text_box(slide, "Tabla de funciones de agregación:",
             Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
agg_funcs = [
    ("np.sum()",    "np.nansum()",    "Suma"),
    ("np.prod()",   "np.nanprod()",   "Producto"),
    ("np.mean()",   "np.nanmean()",   "Media"),
    ("np.std()",    "np.nanstd()",    "Desv. estándar"),
    ("np.var()",    "np.nanvar()",    "Varianza"),
    ("np.min()",    "np.nanmin()",    "Mínimo"),
    ("np.max()",    "np.nanmax()",    "Máximo"),
    ("np.median()", "np.nanmedian()", "Mediana"),
    ("np.argmin()", "np.nanargmin()", "Índice del mínimo"),
    ("np.argmax()", "np.nanargmax()", "Índice del máximo"),
]
add_rect(slide, Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.3), fill_color=NAVY)
for lbl, x in [("Función", 6.7), ("Versión con NaN", 8.75), ("Descripción", 11.0)]:
    add_text_box(slide, lbl, Inches(x), Inches(1.47), Inches(1.9), Inches(0.22),
                 font_size=10, bold=True, color=WHITE)
y = Inches(1.72)
for i, (fn, nan_fn, desc) in enumerate(agg_funcs):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(6.6), y, Inches(6.3), Inches(0.27), fill_color=bg)
    add_text_box(slide, fn,     Inches(6.7),  y + Inches(0.04), Inches(1.9), Inches(0.22), font_size=9.5, bold=True, color=BLUE, font_name="Courier New")
    add_text_box(slide, nan_fn, Inches(8.75), y + Inches(0.04), Inches(2.1), Inches(0.22), font_size=9.5, color=DARK_GRAY, font_name="Courier New")
    add_text_box(slide, desc,   Inches(11.0), y + Inches(0.04), Inches(1.8), Inches(0.22), font_size=9.5, color=DARK_GRAY)
    y += Inches(0.27)


# ══════════════════════════════════════════════════════════════════════════════
# MÓDULO 4
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Módulo 4", "Comparación · Masking · Fancy Indexing · Ordenamiento · Reemplazo")

# SLIDE 19 — Comparación y operadores booleanos
slide = content_slide("Operadores de comparación y lógicos: la base del filtrado vectorizado")
add_code_block(slide,
    "x = np.arange(1, 10)\n\n"
    "x > 3     # [F, F, F, T, T, T, T, T, T]\n"
    "x < 3     # [T, T, F, F, F, F, F, F, F]\n"
    "x == 3    # [F, F, T, F, F, F, F, F, F]\n"
    "x != 3    # [T, T, F, T, T, T, T, T, T]\n\n"
    "# Consultas sobre el arreglo booleano\n"
    "np.sum(M > 5)          # cuántos > 5\n"
    "np.sum(M > 5, axis=1)  # cuántos > 5 por fila\n"
    "np.any(M > 8)          # ¿hay alguno > 8?\n"
    "np.all(M < 10)         # ¿todos < 10?",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(3.1), font_size=11)

add_text_box(slide, "Operadores lógicos elemento a elemento:",
             Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.32),
             font_size=12, bold=True, color=BLUE)
ops = [("&",  "np.bitwise_and()", "AND"),
       ("|",  "np.bitwise_or()",  "OR"),
       ("^",  "np.bitwise_xor()", "XOR"),
       ("~",  "np.bitwise_not()", "NOT")]
add_rect(slide, Inches(6.6), Inches(1.42), Inches(6.3), Inches(0.3), fill_color=NAVY)
for lbl, x in [("Símbolo", 6.7), ("ufunc equivalente", 7.6), ("Operación", 10.0)]:
    add_text_box(slide, lbl, Inches(x), Inches(1.47), Inches(2.2), Inches(0.22),
                 font_size=11, bold=True, color=WHITE)
y = Inches(1.72)
for i, (sym, fn, op) in enumerate(ops):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(6.6), y, Inches(6.3), Inches(0.32), fill_color=bg)
    add_text_box(slide, sym, Inches(6.7),  y + Inches(0.05), Inches(0.7), Inches(0.26), font_size=14, bold=True, color=RED, font_name="Courier New")
    add_text_box(slide, fn,  Inches(7.6),  y + Inches(0.05), Inches(2.2), Inches(0.26), font_size=10, color=DARK_GRAY, font_name="Courier New")
    add_text_box(slide, op,  Inches(10.0), y + Inches(0.05), Inches(2.7), Inches(0.26), font_size=11, color=DARK_GRAY)
    y += Inches(0.32)

add_code_block(slide,
    "# Histórico chancador primario: µ=4850 tph, σ=750\n"
    "hist = rng.normal(1250, 750, 100) + 3600\n\n"
    "np.sum(hist < 4000)                      # días < 4000 tph\n"
    "np.sum((hist >= 4500) & (hist <= 5000))  # entre 4500 y 5000\n"
    "np.sum((hist > 6000) | (hist < 4000))    # rendimiento extremo",
    Inches(6.6), Inches(2.85), Inches(6.3), Inches(2.0), font_size=10.5)

add_callout(slide, "¡Atención!",
            "Los paréntesis son OBLIGATORIOS al combinar condiciones. "
            "La precedencia de & y | difiere de and/or.",
            Inches(0.4), Inches(4.3), Inches(5.9), Inches(0.75),
            accent=RGBColor(0xE6, 0x7E, 0x22), bg=RGBColor(0xFE, 0xF9, 0xF0))

# SLIDE 20 — Masking y Fancy Indexing
slide = content_slide("Masking y Fancy Indexing: filtrado y selección vectorizados")
add_callout(slide, "Masking",
            "Usar un arreglo booleano como índice para seleccionar elementos que cumplen una condición. "
            "Es la herramienta central del preprocesamiento vectorizado.",
            Inches(0.4), Inches(1.05), Inches(5.9), Inches(0.85))

add_code_block(slide,
    "# Masking básico\n"
    "M[M < 4]              # todos los menores que 4\n"
    "M[(M > 3) & (M < 8)]  # compuesta\n"
    "# resultado: siempre arreglo 1D\n\n"
    "# Campaña metalúrgica: filtrar muestras\n"
    "# P cols: [tph, psi, %malla, %Cu]\n"
    "mask = (\n"
    "    (P[:, 0] >= 5000) &   # tph ≥ 5000\n"
    "    (P[:, 1] >  11.0) &   # psi > 11\n"
    "    (P[:, 3] >= 0.62)     # Cu% ≥ 0.62\n"
    ")\n"
    "P[mask]             # filas que cumplen todo\n"
    "P[mask][:, [0, 3]]  # + selección de columnas",
    Inches(0.4), Inches(2.05), Inches(6.0), Inches(3.5), font_size=10.5)

add_callout(slide, "Fancy Indexing",
            "Pasar un arreglo de índices para seleccionar múltiples posiciones a la vez. "
            "La geometría del resultado = geometría de los índices.",
            Inches(6.6), Inches(1.05), Inches(6.3), Inches(0.85))

add_code_block(slide,
    "x = rng.integers(1, 100, size=20)\n\n"
    "# Selección por lista de índices\n"
    "x[[13, 7, 18]]       # → arreglo (3,)\n\n"
    "# Índices en forma de matriz → mismo shape\n"
    "idx = np.array([[13, 7], [18, 2]])\n"
    "x[idx]               # → arreglo (2, 2)\n\n"
    "# En 2D: selección pareada vs rectangular\n"
    "row = np.array([0, 2, 4])\n"
    "col = np.array([1, 3, 5])\n"
    "M[row, col]           # (3,) — pareado\n"
    "M[np.ix_(row, col)]   # (3,3) — rectangular",
    Inches(6.6), Inches(2.05), Inches(6.3), Inches(3.3), font_size=10.5)

# SLIDE 21 — Ordenamiento y reemplazo
slide = content_slide("Ordenamiento, particionamiento y reemplazo condicional")
add_code_block(slide,
    "u = rng.integers(1, 20, size=10)\n\n"
    "np.sort(u)          # copia ordenada (no modifica u)\n"
    "u.sort()            # modifica u en sitio\n\n"
    "# Índices que ordenarían el arreglo\n"
    "idx = np.argsort(u)\n"
    "u[idx]              # equivale a np.sort(u)\n\n"
    "# Ordenamiento por eje en 2D\n"
    "np.sort(M, axis=0)  # ordena cada columna\n"
    "np.sort(M, axis=1)  # ordena cada fila",
    Inches(0.4), Inches(1.05), Inches(6.0), Inches(2.7), font_size=11)

add_code_block(slide,
    "# Particionamiento: K menores al inicio\n"
    "np.partition(u, 3)     # los 3 primeros ← mínimos\n"
    "np.partition(M, 2, axis=0)\n\n"
    "# K vecinos más cercanos\n"
    "dist = np.zeros((30, 30))  # distancias euclidianas\n"
    "for i in range(30):\n"
    "    for j in range(30):\n"
    "        dist[i,j] = np.sum((P[i]-P[j])**2)\n\n"
    "K = 3\n"
    "nearest = np.argpartition(dist, K+1, axis=1)\n"
    "nearest[:, 1:K+1]  # excluye el punto mismo",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(2.95), font_size=10.5)

add_code_block(slide,
    "# Reemplazo condicional: np.where(cond, si_true, si_false)\n"
    "T = rng.normal(1800, 650, 1000)   # tph molino\n\n"
    "# Política: tratar < 1200 tph como dato ausente\n"
    "T_corr = np.where(T < 1200, np.nan, T)\n\n"
    "# Recorte en rango\n"
    "T_clip = np.clip(T, 1200, 5000)   # recorta entre límites",
    Inches(0.4), Inches(3.9), Inches(12.5), Inches(1.85), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# MÓDULO 5
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Módulo 5", "NumPy como base del modelamiento predictivo")

# SLIDE 22 — Regresión lineal: modelo y costo
slide = content_slide("Regresión lineal: el modelo y su función de costo (MSE)")
add_text_box(slide, "El modelo:",
             Inches(0.4), Inches(1.05), Inches(5.9), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
add_text_box(slide,
    "ŷᵢ = w₀ + Σ wⱼ xᵢⱼ  =  Xw + b",
    Inches(0.4), Inches(1.42), Inches(5.9), Inches(0.4),
    font_size=15, bold=True, color=NAVY, font_name="Courier New")
add_text_box(slide,
    "X ∈ ℝᵐˣⁿ  (m obs., n vars.)    w ∈ ℝⁿ  (parámetros)    b = w₀  (sesgo)",
    Inches(0.4), Inches(1.88), Inches(5.9), Inches(0.32),
    font_size=11, color=DARK_GRAY, font_name="Courier New")

add_text_box(slide, "Función de costo — Error Cuadrático Medio (MSE):",
             Inches(0.4), Inches(2.35), Inches(5.9), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
add_text_box(slide,
    "MSE = (1/m) Σᵢ (yᵢ − ŷᵢ)²",
    Inches(0.4), Inches(2.72), Inches(5.9), Inches(0.38),
    font_size=15, bold=True, color=NAVY, font_name="Courier New")

add_text_box(slide, "Gradiente descendente (GD):",
             Inches(0.4), Inches(3.2), Inches(5.9), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
add_text_box(slide,
    "∂MSE/∂w = -(2/m) Xᵀ(y − ŷ)\n"
    "∂MSE/∂b = -(2/m) Σ(yᵢ − ŷᵢ)\n\n"
    "wₖ₊₁ = wₖ − α · ∂MSE/∂w\n"
    "bₖ₊₁ = bₖ − α · ∂MSE/∂b",
    Inches(0.4), Inches(3.57), Inches(5.9), Inches(1.5),
    font_size=12, color=DARK_GRAY, font_name="Courier New")

add_code_block(slide,
    "def fit_linear_regression(\n"
    "    X, y, iterations=100,\n"
    "    tolerance=1e-5, learning_rate=0.5\n"
    "):\n"
    "    y = y.reshape(-1, 1)\n"
    "    # sesgo como primera columna de 1s\n"
    "    X = np.concatenate(\n"
    "        [np.ones((X.shape[0], 1)), X], axis=1\n"
    "    )\n"
    "    w = rng.random(size=(X.shape[1], 1))\n\n"
    "    for _ in range(iterations):\n"
    "        y_pred   = X @ w\n"
    "        gradient = 2 * X.T @ (y_pred - y) / X.shape[0]\n"
    "        w       -= gradient * learning_rate\n"
    "        if np.linalg.norm(gradient) < tolerance:\n"
    "            break\n"
    "    return w",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(5.55), font_size=10.5)

# SLIDE 23 — R² y advertencias
slide = content_slide("Evaluación de la regresión lineal: el coeficiente R²")
add_text_box(slide,
    "r² = 1 − SSres/SStot      SSres = Σ(yᵢ − ŷᵢ)²      SStot = Σ(yᵢ − ȳ)²",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.42),
    font_size=14, bold=True, color=NAVY, font_name="Courier New")
add_text_box(slide,
    "Interpretación: fracción de la variabilidad de y explicada por el modelo. "
    "r² = 1 → ajuste perfecto. r² = 0 → el modelo no explica nada.",
    Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4),
    font_size=12, color=DARK_GRAY)

add_code_block(slide,
    "def squared_r(y_real, y_pred):\n"
    "    SS_res = np.sum((y_real - y_pred)**2)\n"
    "    SS_tot = np.sum((y_real - np.mean(y_real))**2)\n"
    "    return 1 - SS_res / SS_tot\n\n"
    "y_train_pred = X_train @ w[1:] + w[0]\n"
    "y_test_pred  = X_test  @ w[1:] + w[0]\n\n"
    "r2_train = squared_r(y_train, y_train_pred)\n"
    "r2_test  = squared_r(y_test,  y_test_pred)",
    Inches(0.4), Inches(2.1), Inches(6.0), Inches(2.5), font_size=12)

advertencias = [
    ("Escalar ANTES de ajustar.",
     "Variables en distintas unidades (tph, psi, %) sesgan el gradiente. Aplicar z-score."),
    ("Outliers distorsionan el MSE.",
     "El MSE es cuadrático: un valor extremo tiene peso desproporcionado."),
    ("r² alto ≠ modelo útil.",
     "Puede haber sobreajuste, autocorrelación en residuos, o relación no lineal."),
    ("Los datos reales tienen NaN.",
     "Leyes extremas, registros duplicados. La limpieza no es opcional."),
]
y = Inches(2.1)
for i, (titulo, desc) in enumerate(advertencias):
    add_rect(slide, Inches(6.6), y, Inches(0.35), Inches(0.7), fill_color=RED if i == 0 else NAVY)
    add_text_box(slide, str(i+1), Inches(6.6), y + Inches(0.17), Inches(0.35), Inches(0.32),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, titulo, Inches(7.1), y, Inches(5.8), Inches(0.28),
                 font_size=11, bold=True, color=BLUE)
    add_text_box(slide, desc, Inches(7.1), y + Inches(0.3), Inches(5.8), Inches(0.4),
                 font_size=10.5, color=DARK_GRAY)
    y += Inches(0.82)

# SLIDE 24 — Regresión logística
slide = content_slide("Regresión logística binaria: cuando la respuesta es 0 o 1")
add_text_box(slide, "El problema de clasificación binaria:",
             Inches(0.4), Inches(1.05), Inches(5.9), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
add_text_box(slide,
    "La variable de respuesta es binaria (y ∈ {0, 1}). "
    "Un modelo lineal puede predecir valores fuera de [0,1]: no tiene sentido como probabilidad.\n\n"
    "Solución: transformar la salida con la función logística (sigmoide):",
    Inches(0.4), Inches(1.42), Inches(5.9), Inches(0.8),
    font_size=11, color=DARK_GRAY)
add_text_box(slide,
    "φ(u) = 1 / (1 + exp(−u))",
    Inches(0.4), Inches(2.27), Inches(5.9), Inches(0.38),
    font_size=16, bold=True, color=NAVY, font_name="Courier New")
add_text_box(slide,
    "φ: ℝ → (0,1)  ·  Interpretación: P(y=1 | x)",
    Inches(0.4), Inches(2.72), Inches(5.9), Inches(0.32),
    font_size=11, color=DARK_GRAY, font_name="Courier New")

add_text_box(slide, "Función de costo — Entropía cruzada binaria:",
             Inches(0.4), Inches(3.15), Inches(5.9), Inches(0.32),
             font_size=13, bold=True, color=BLUE)
add_text_box(slide,
    "ℒ = −(1/m) Σ [yᵢ log(p̂ᵢ) + (1−yᵢ) log(1−p̂ᵢ)]",
    Inches(0.4), Inches(3.52), Inches(5.9), Inches(0.38),
    font_size=13, bold=True, color=NAVY, font_name="Courier New")

add_text_box(slide, "Propiedad clave: φ'(u) = φ(u)(1 − φ(u))",
             Inches(0.4), Inches(4.0), Inches(5.9), Inches(0.32),
             font_size=11, italic=True, color=DARK_GRAY)

add_code_block(slide,
    "def sigmoid(x):\n"
    "    return 1 / (1 + np.exp(-x))\n\n"
    "def log_loss(y, y_pred):\n"
    "    eps    = 1e-12\n"
    "    y_pred = np.clip(y_pred, eps, 1 - eps)\n"
    "    return -np.mean(\n"
    "        y * np.log(y_pred) +\n"
    "        (1 - y) * np.log(1 - y_pred)\n"
    "    )\n\n"
    "def fit_logistic_regression(X, y, ...):\n"
    "    X = np.hstack((np.ones((X.shape[0],1)), X))\n"
    "    w = np.zeros(X.shape[1])\n"
    "    for _ in range(iterations):\n"
    "        z        = np.dot(X, w)\n"
    "        y_pred   = sigmoid(z)\n"
    "        gradient = np.dot(X.T, y_pred-y) / y.size\n"
    "        w       -= learning_rate * gradient\n"
    "    return w",
    Inches(6.6), Inches(1.05), Inches(6.3), Inches(5.6), font_size=10.5)

# SLIDE 25 — Evaluación clasificación
slide = content_slide("Evaluación del clasificador: matriz de confusión y métricas")
add_text_box(slide,
    "Umbral de clasificación: predecir Ŷ=1 si P(Y=1) > 0.5",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35),
    font_size=13, bold=True, color=BLUE)

# Matriz de confusión visual
tabla_cm = [
    ("",           "Real Y = 1",        "Real Y = 0"),
    ("Pred Ŷ = 1", "VP (Verd. Pos.)",   "FP (Falso Pos.)"),
    ("Pred Ŷ = 0", "FN (Falso Neg.)",   "VN (Verd. Neg.)"),
]
y = Inches(1.5)
for i, row in enumerate(tabla_cm):
    for j, cell in enumerate(row):
        x = Inches(0.4 + j * 2.0)
        bg = NAVY if i == 0 or j == 0 else (RGBColor(0xE8,0xF5,0xE9) if (i==1 and j==1) or (i==2 and j==2) else RGBColor(0xFF,0xEB,0xEE))
        add_rect(slide, x, y, Inches(1.95), Inches(0.5), fill_color=bg)
        fc = WHITE if (i == 0 or j == 0) else DARK_GRAY
        add_text_box(slide, cell, x + Inches(0.05), y + Inches(0.1), Inches(1.85), Inches(0.35),
                     font_size=11, bold=(i == 0 or j == 0), color=fc, align=PP_ALIGN.CENTER)
    y += Inches(0.5)

add_text_box(slide, "Sensibilidad (Recall) = VP / (VP + FN)  →  Fracción de positivos reales bien detectados",
             Inches(0.4), Inches(3.0), Inches(5.9), Inches(0.35),
             font_size=11, color=DARK_GRAY)
add_text_box(slide, "Especificidad         = VN / (VN + FP)  →  Fracción de negativos reales bien detectados",
             Inches(0.4), Inches(3.4), Inches(5.9), Inches(0.35),
             font_size=11, color=DARK_GRAY)

add_code_block(slide,
    "def confusion_matrix(y_true, y_pred):\n"
    "    VP = np.sum((y_true==1) & (y_pred==1))\n"
    "    VN = np.sum((y_true==0) & (y_pred==0))\n"
    "    FP = np.sum((y_true==0) & (y_pred==1))\n"
    "    FN = np.sum((y_true==1) & (y_pred==0))\n"
    "    return np.array([[VP, FP], [FN, VN]])\n\n"
    "def predict(X, theta):\n"
    "    X      = np.hstack((np.ones((X.shape[0],1)), X))\n"
    "    y_pred = sigmoid(np.dot(X, theta))\n"
    "    return np.where(y_pred > 0.5, 1, 0)\n\n"
    "cm = confusion_matrix(y_test, predict(X_test, w))\n"
    "sensibilidad  = cm[0,0] / (cm[0,0] + cm[1,0])\n"
    "especificidad = cm[1,1] / (cm[1,1] + cm[0,1])",
    Inches(6.6), Inches(1.42), Inches(6.3), Inches(3.7), font_size=10.5)

add_callout(slide, "Detección de fallas en minería",
            "Si el modelo detecta fallas en equipos (y=1 = falla), maximizar sensibilidad es crítico. "
            "Un falso negativo (falla no detectada) es peor que un falso positivo (parada innecesaria).",
            Inches(0.4), Inches(3.9), Inches(5.9), Inches(1.0),
            accent=RGBColor(0x1D, 0x9A, 0x6C), bg=RGBColor(0xF0, 0xFA, 0xF5))


# ══════════════════════════════════════════════════════════════════════════════
# CIERRE
# ══════════════════════════════════════════════════════════════════════════════
section_divider("Cierre", "NumPy como base de todo lo que viene")

# SLIDE 26 — Resumen
slide = content_slide("Hoja de referencia rápida — NumPy en una sola página")

col1 = [
    ("CREACIÓN", True, False, 11, BLUE),
    ("np.array([1,2,3])", False, False, 10, CODE_FG),
    ("np.zeros/ones/full/eye", False, False, 10, CODE_FG),
    ("np.arange / np.linspace", False, False, 10, CODE_FG),
    ("np.meshgrid(x, y)", False, False, 10, CODE_FG),
    ("rng.normal / integers", False, False, 10, CODE_FG),
    ("", False, False, 8, WHITE),
    ("ATRIBUTOS", True, False, 11, BLUE),
    (".shape  .ndim  .size", False, False, 10, CODE_FG),
    (".dtype  .itemsize  .nbytes", False, False, 10, CODE_FG),
    ("", False, False, 8, WHITE),
    ("RESHAPE / UNIÓN", True, False, 11, BLUE),
    ("A.T     A.ravel()   A.reshape()", False, False, 10, CODE_FG),
    ("np.vstack / hstack / concatenate", False, False, 10, CODE_FG),
    ("np.vsplit / hsplit / split", False, False, 10, CODE_FG),
    ("", False, False, 8, WHITE),
    ("I/O", True, False, 11, BLUE),
    ("np.loadtxt / genfromtxt", False, False, 10, CODE_FG),
    ("np.save / load   (.npy)", False, False, 10, CODE_FG),
]

col2 = [
    ("OPERATORIA", True, False, 11, BLUE),
    ("a+b  a-b  a*b  a/b  a**2", False, False, 10, CODE_FG),
    ("A @ B          (matricial)", False, False, 10, CODE_FG),
    ("np.dot(a, b)   (producto interno)", False, False, 10, CODE_FG),
    ("np.linalg.norm / solve / svd", False, False, 10, CODE_FG),
    ("np.abs / exp / log / sin / cos", False, False, 10, CODE_FG),
    ("", False, False, 8, WHITE),
    ("AGREGACIÓN", True, False, 11, BLUE),
    ("A.sum/mean/std/var(axis=0)", False, False, 10, CODE_FG),
    ("A.max/min/argmax/argmin()", False, False, 10, CODE_FG),
    ("np.nanmean / nanstd / nansum", False, False, 10, CODE_FG),
    ("", False, False, 8, WHITE),
    ("FILTRADO Y SELECCIÓN", True, False, 11, BLUE),
    ("x > 3   np.any/all(x > 3)", False, False, 10, CODE_FG),
    ("x[(x>3)&(x<8)]  (masking)", False, False, 10, CODE_FG),
    ("x[[0,4,7]]      (fancy idx)", False, False, 10, CODE_FG),
    ("M[np.ix_(rows, cols)]", False, False, 10, CODE_FG),
    ("np.where(cond, a, b)", False, False, 10, CODE_FG),
    ("np.sort / argsort / partition", False, False, 10, CODE_FG),
]

for col_data, x_start in [(col1, 0.4), (col2, 6.8)]:
    txBox = slide.shapes.add_textbox(Inches(x_start), Inches(1.05),
                                     Inches(6.0), Inches(6.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in col_data:
        text, bold, italic, sz, col_c = item
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.name = "Courier New" if not bold else "Calibri"
        run.font.color.rgb = col_c

# SLIDE 27 — Resumen modulos
slide = content_slide("Lo que aprendimos: del arreglo al modelo predictivo")
modulos_resumen = [
    ("Módulo 1", "El arreglo",          "ndarray, geometría, ejes, atributos, tipos, rutinas de creación, I/O, indexación, slicing, reshape"),
    ("Módulo 2", "Operatoria",           "ufuncs, aritmética vectorizada, Hadamard vs producto matricial, álgebra lineal, trigonometría, log/exp"),
    ("Módulo 3", "Broadcasting",         "Reglas de broadcasting, z-score por broadcasting, funciones de agregación por eje"),
    ("Módulo 4", "Filtrado",             "Comparación booleana, operadores lógicos, masking, fancy indexing, ordenamiento, reemplazo condicional"),
    ("Módulo 5", "Modelamiento",         "Regresión lineal (MSE + GD), R², regresión logística binaria (sigmoide + log-loss), matriz de confusión"),
]
y = Inches(1.1)
for i, (mod, nombre, desc) in enumerate(modulos_resumen):
    c = [NAVY, BLUE, RGBColor(0x1D, 0x9A, 0x6C), RED, RGBColor(0xE6, 0x7E, 0x22)][i]
    add_rect(slide, Inches(0.4), y, Inches(1.6), Inches(0.9), fill_color=c)
    add_text_box(slide, mod,    Inches(0.42), y + Inches(0.06), Inches(1.55), Inches(0.3),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, nombre, Inches(0.42), y + Inches(0.38), Inches(1.55), Inches(0.3),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, Inches(2.15), y + Inches(0.2), Inches(11.0), Inches(0.55),
                 font_size=11, color=DARK_GRAY)
    y += Inches(1.05)

add_text_box(slide,
    "Dominar NumPy no es dominar una librería. Es dominar la lógica sobre la cual "
    "todo el stack científico de Python está construido.",
    Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.55),
    font_size=12, italic=True, bold=True, color=NAVY)

# SLIDE FINAL
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=NAVY)
add_rect(slide, Inches(0.6), Inches(3.2), Inches(11.8), Pt(3), fill_color=RED)
add_text_box(slide, "Gracias",
             Inches(0.6), Inches(3.35), Inches(11.8), Inches(0.9),
             font_size=48, bold=True, color=WHITE, font_name="Calibri")
add_text_box(slide, "Notas completas disponibles en el blog del curso",
             Inches(0.6), Inches(4.35), Inches(11.8), Inches(0.4),
             font_size=16, color=MID_GRAY)
add_text_box(slide, "Próxima sesión: Análisis tabular con Pandas",
             Inches(0.6), Inches(4.85), Inches(11.8), Inches(0.35),
             font_size=14, color=RED)


# ── Guardar ───────────────────────────────────────────────────────────────────
output_path = "NumPy_Presentacion_ICM.pptx"
prs.save(output_path)
print(f"✓ Presentación guardada: {output_path}")
print(f"  Total de slides: {len(prs.slides)}")
